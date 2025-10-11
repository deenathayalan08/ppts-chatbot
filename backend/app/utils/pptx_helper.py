from pptx import Presentation
from pptx.util import Inches, Pt
import os
from app.models.ai_model import generate_slide_contents  # AI slide content generator

# --- Apply theme ---
def apply_theme(slide, theme):
    theme_colors = {
        "classic": ("000000", "FFFFFF"),
        "modern": ("1F4E79", "DDEBF7"),
        "professional": ("385723", "E2F0D9"),
        "minimal": ("FFFFFF", "333333")
    }
    title_color, body_color = theme_colors.get(theme, ("000000", "FFFFFF"))

    if slide.shapes.title:
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(20)


# --- Generate PPT using AI model ---
def generate_ppt_from_topic(topic, num_slides=5, theme="classic",
                            audience="Beginner", purpose="Informative",
                            tone="Formal", auto_media=False, save_path=None):
    slides_data = []
    prs = Presentation()

    try:
        # Use AI model to generate titles/body
        ai_slides = generate_slide_contents(topic, num_slides)

        for i, slide_info in enumerate(ai_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_text = slide_info["title"]
            body_text = slide_info["body"]

            # Apply theme
            slide.shapes.title.text = title_text
            slide.placeholders[1].text = body_text
            apply_theme(slide, theme)

            # Auto media placeholder
            media = None
            if auto_media:
                media = "Image/Chart Placeholder"

            slides_data.append({
                "slide_number": i + 1,
                "title": title_text,
                "body": body_text,
                "theme": theme,
                "media_suggestion": media
            })

        if save_path:
            prs.save(save_path)

        return slides_data

    except Exception as e:
        print(f"[pptx_helper] Error generating PPT: {e}")
        raise


# --- Update slide ---
def update_slide(filename, slide_number, title, text, folder="generated_files"):
    file_path = os.path.join(folder, filename)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"{file_path} does not exist!")

    prs = Presentation(file_path)
    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError("Invalid slide number.")

    slide = prs.slides[slide_number - 1]
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = text
    prs.save(file_path)
    return file_path


# --- Add image ---
def add_image_to_slide(filename, slide_number, image_path, folder="generated_files"):
    file_path = os.path.join(folder, filename)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"{file_path} does not exist!")

    prs = Presentation(file_path)
    slide = prs.slides[slide_number - 1]
    slide.shapes.add_picture(image_path, Inches(1), Inches(2), height=Inches(3))
    prs.save(file_path)
    return file_path


# --- Add video ---
def add_video_to_slide(filename, slide_number, video_path, folder="generated_files"):
    file_path = os.path.join(folder, filename)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"{file_path} does not exist!")

    prs = Presentation(file_path)
    slide = prs.slides[slide_number - 1]
    slide.shapes.add_movie(video_path, Inches(1), Inches(1), width=Inches(6), height=Inches(3))
    prs.save(file_path)
    return file_path
