from fastapi import APIRouter, HTTPException, Query
from fastapi.responses import FileResponse, JSONResponse
from pptx import Presentation
import os
import logging
from app.config import UPLOAD_FOLDER
from app.models.ai_model import generate_slide_contents

router = APIRouter()

@router.post("/generate/")
async def generate_ppt(topic: str, preview: bool = Query(True)):
    try:
        slides_data = generate_slide_contents(topic)

        if preview:
            # Return JSON for frontend preview
            return JSONResponse(content={"slides": slides_data})

        # Generate PPT file
        prs = Presentation()
        for slide_content in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_content["title"]
            slide.placeholders[1].text = slide_content["body"]

        # Safe filename
        safe_topic = "".join(c if c.isalnum() else "_" for c in topic).lower()
        filename = f"{safe_topic}.pptx"
        filepath = os.path.join(UPLOAD_FOLDER, filename)

        prs.save(filepath)
        logging.info(f"✅ PPT created for topic: {topic}")

        # Ensure proper download headers
        return FileResponse(
            path=filepath,
            filename=filename,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        logging.error(f"❌ Error generating PPT: {e}")
        raise HTTPException(status_code=500, detail=str(e))
