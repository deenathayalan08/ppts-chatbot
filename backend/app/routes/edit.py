from fastapi import APIRouter, UploadFile, Form
from fastapi.responses import JSONResponse
from pptx import Presentation
import os

router = APIRouter()

UPLOAD_FOLDER = "uploads"

# --- Preview slides ---
@router.get("/preview/")
async def preview(filename: str):
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    if not os.path.exists(file_path):
        return JSONResponse({"error": "File not found"}, status_code=404)

    prs = Presentation(file_path)
    slides_data = []

    for slide in prs.slides:
        title = slide.shapes.title.text if slide.shapes.title else ""
        body = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                body += shape.text + "\n"
        slides_data.append({"title": title, "body": body.strip()})

    return {"slides": slides_data, "filename": filename}


# --- Update a slide ---
@router.post("/update/")
async def update_slide(slide_number: int = Form(...), title: str = Form(...), text: str = Form(...), filename: str = Form(...)):
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    if not os.path.exists(file_path):
        return JSONResponse({"error": "File not found"}, status_code=404)

    prs = Presentation(file_path)

    if slide_number < 1 or slide_number > len(prs.slides):
        return JSONResponse({"error": "Invalid slide number"}, status_code=400)

    slide = prs.slides[slide_number - 1]

    # Update slide title and body
    if slide.shapes.title:
        slide.shapes.title.text = title
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            shape.text = text

    prs.save(file_path)

    return {"message": f"Slide {slide_number} updated successfully!"}
